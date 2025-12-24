#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Powder Inspection Management System - Presentation Generator
Generates a PowerPoint presentation with system overview, architecture, and specifications
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create PowerPoint presentation for Powder Inspection Management System"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # Slide 1: Title Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(103, 126, 234)  # Purple gradient

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Powder Inspection Management System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Quality Control & Traceability Solution"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "December 2024"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER

    # ==========================================
    # Slide 2: System Overview
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_slide_title(slide, "System Overview")

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Purpose
    p = tf.paragraphs[0]
    p.text = "Purpose"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(tf, "Comprehensive quality control system for powder inspection and blending operations", 18)
    add_bullet(tf, "Complete traceability from raw materials to finished products", 18)
    add_bullet(tf, "Automated LOT tracking using barcode/QR code technology", 18)

    # Key Features
    p = tf.add_paragraph()
    p.text = "Key Features"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(24)
    p.space_after = Pt(12)

    add_bullet(tf, "Web-based application (accessible from any PC)", 18)
    add_bullet(tf, "Real-time data synchronization", 18)
    add_bullet(tf, "Barcode label printing and scanning", 18)
    add_bullet(tf, "Quality inspection management", 18)
    add_bullet(tf, "Complete backward/forward traceability", 18)

    # ==========================================
    # Slide 3: System Architecture
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "System Architecture")

    # Architecture diagram (text-based)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Server
    p = tf.paragraphs[0]
    add_architecture_box(slide, 4, 2.2, 2, 1.2, "Server PC (A)",
                         "• Flask Web Application\n• SQLite Database\n• Central Data Management")

    # Network
    network_box = slide.shapes.add_textbox(Inches(4.5), Inches(3.7), Inches(1), Inches(0.4))
    nf = network_box.text_frame
    nf.text = "Network"
    nf.paragraphs[0].font.size = Pt(14)
    nf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Client 1
    add_architecture_box(slide, 1, 4.5, 2, 1.2, "Client PC 1 (B)",
                         "• Inspection Station\n• Barcode Scanner\n• Label Printer")

    # Client 2
    add_architecture_box(slide, 7, 4.5, 2, 1.2, "Client PC 2 (B)",
                         "• Blending Station\n• Barcode Scanner\n• Label Printer\n• Touch Monitor")

    # ==========================================
    # Slide 4: Main Functions
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Main Functions")

    # Two columns
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(5))
    left_tf = left_box.text_frame

    # Left column
    p = left_tf.paragraphs[0]
    p.text = "1. Incoming Powder Inspection"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    add_bullet(left_tf, "Quality testing of raw materials", 14)
    add_bullet(left_tf, "LOT number registration", 14)
    add_bullet(left_tf, "Multi-criteria inspection results", 14)

    p = left_tf.add_paragraph()
    p.text = "2. Blending Work Planning"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(left_tf, "Recipe management", 14)
    add_bullet(left_tf, "Work order creation", 14)
    add_bullet(left_tf, "Production scheduling", 14)

    p = left_tf.add_paragraph()
    p.text = "3. Material Input Tracking"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(left_tf, "Barcode scanning for LOT input", 14)
    add_bullet(left_tf, "Weight measurement recording", 14)
    add_bullet(left_tf, "Cross-contamination detection", 14)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
    right_tf = right_box.text_frame

    p = right_tf.paragraphs[0]
    p.text = "4. Quality Inspection Results"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    add_bullet(right_tf, "Comprehensive test result storage", 14)
    add_bullet(right_tf, "Pass/Fail determination", 14)
    add_bullet(right_tf, "Historical data management", 14)

    p = right_tf.add_paragraph()
    p.text = "5. Complete Traceability"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(right_tf, "Backward: Finished → Raw materials", 14)
    add_bullet(right_tf, "Forward: Raw material → Products", 14)
    add_bullet(right_tf, "Complete audit trail", 14)

    p = right_tf.add_paragraph()
    p.text = "6. Barcode/Label System"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(right_tf, "Automatic label generation", 14)
    add_bullet(right_tf, "QR/Barcode printing", 14)
    add_bullet(right_tf, "Wireless scanner integration", 14)

    # ==========================================
    # Slide 5: User Interface - Menu Structure
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "User Interface - Menu Structure")

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = content_box.text_frame

    menu_items = [
        "1. Dashboard - System overview and ongoing inspections",
        "2. Incoming Powder Inspection - Raw material quality testing",
        "3. Blending Work Planning - Recipe and work order management",
        "4. Blending Work - Material input and production execution",
        "5. Blending Work Status - Real-time production monitoring",
        "6. Blending Powder Inspection - Finished product quality testing",
        "7. Inspection Results - Historical data query",
        "8. REWORK - Quality issue handling (future)",
        "9. Traceability - LOT tracking and audit trail",
        "10. Admin Mode - System configuration and user management"
    ]

    for item in menu_items:
        add_bullet(tf, item, 16)

    # ==========================================
    # Slide 6: Workflow
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Operational Workflow")

    # Workflow steps
    steps = [
        ("1. Incoming\nInspection", 1, 2.2),
        ("2. Recipe\nSetup", 3, 2.2),
        ("3. Work Order\nCreation", 5, 2.2),
        ("4. Material\nInput", 7, 2.2),
        ("5. Blending\nExecution", 1, 4.5),
        ("6. Label\nPrinting", 3, 4.5),
        ("7. Quality\nInspection", 5, 4.5),
        ("8. Traceability\nTracking", 7, 4.5)
    ]

    for step_text, x, y in steps:
        add_workflow_box(slide, x, y, 1.5, 1, step_text)

    # Add arrows
    add_arrow(slide, 2.6, 2.7, 2.9, 2.7)  # 1 -> 2
    add_arrow(slide, 4.6, 2.7, 4.9, 2.7)  # 2 -> 3
    add_arrow(slide, 6.6, 2.7, 6.9, 2.7)  # 3 -> 4
    add_arrow(slide, 7.75, 3.3, 7.75, 4.4)  # 4 -> down
    add_arrow(slide, 7.25, 5, 4.6, 5)  # down -> 7
    add_arrow(slide, 4.1, 5, 2.6, 5)  # 7 -> 6
    add_arrow(slide, 2.1, 5, 1.75, 5)  # 6 -> 5
    add_arrow(slide, 1.2, 5.5, 1.2, 6.5)  # 5 -> down
    add_arrow(slide, 1.7, 6.7, 6.9, 6.7)  # across to 8

    # ==========================================
    # Slide 7: Hardware Configuration
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Hardware Configuration")

    # Create table
    rows, cols = 8, 4
    left = Inches(0.8)
    top = Inches(2)
    width = Inches(8.4)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.5)

    # Header row
    headers = ['Equipment', 'Qty', 'Model', 'Purpose']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    equipment_data = [
        ['Desktop PC (Server)', '1', 'SBD71223', 'Central server & database'],
        ['Fanless Industrial PC', '2', 'HDL-BOX90C', 'Inspection & blending stations'],
        ['Monitor 24"', '1', 'LG 24MR400W', 'Server display'],
        ['Touch Monitor 24"', '1', 'LG 24MR400', 'Blending station'],
        ['Label Printer', '2', 'BIXOLON SLP-TX420/400', 'Barcode label printing'],
        ['Barcode Scanner', '2', 'Zebra DS2278', 'LOT number scanning'],
        ['Total Investment', '9 units', '', 'KRW 5,145,000']
    ]

    for row_idx, row_data in enumerate(equipment_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if row_idx == 7:  # Total row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
                cell.text_frame.paragraphs[0].font.bold = True

    # ==========================================
    # Slide 8: Barcode & Traceability System
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Barcode & Traceability System")

    # Left: Barcode Features
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.2), Inches(4.5))
    left_tf = left_box.text_frame

    p = left_tf.paragraphs[0]
    p.text = "Barcode Label System"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(left_tf, "Format: CODE128 / QR Code", 14)
    add_bullet(left_tf, "Content: Product name, LOT number, Pack number, Weight, Date", 14)
    add_bullet(left_tf, "Size: 100mm × 100mm thermal labels", 14)
    add_bullet(left_tf, "Printer: BIXOLON thermal transfer", 14)
    add_bullet(left_tf, "Scanner: Zebra DS2278 wireless 2D", 14)

    p = left_tf.add_paragraph()
    p.text = "\nBarcode Example:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(18)

    p = left_tf.add_paragraph()
    p.text = "PN:High-Strength-Mix|LOT:20241224-001|1/5"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'

    # Right: Traceability Features
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.2), Inches(4.5))
    right_tf = right_box.text_frame

    p = right_tf.paragraphs[0]
    p.text = "Traceability Features"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(right_tf, "Backward Traceability:\nFinished product → Raw materials", 14)
    add_bullet(right_tf, "Forward Traceability:\nRaw material → All products using it", 14)
    add_bullet(right_tf, "Complete audit trail with timestamps", 14)
    add_bullet(right_tf, "Quality inspection results linked", 14)
    add_bullet(right_tf, "Operator and equipment tracking", 14)

    p = right_tf.add_paragraph()
    p.text = "\nBenefits:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(18)

    add_bullet(right_tf, "Rapid quality issue identification", 14)
    add_bullet(right_tf, "Regulatory compliance", 14)
    add_bullet(right_tf, "Recall capability", 14)

    # ==========================================
    # Slide 9: Expected Benefits
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Expected Benefits")

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame

    # Quality Improvement
    p = tf.paragraphs[0]
    p.text = "🎯 Quality Improvement"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(10)

    add_bullet(tf, "Systematic quality inspection process", 16)
    add_bullet(tf, "Digital record keeping eliminates paper-based errors", 16)
    add_bullet(tf, "Real-time quality monitoring", 16)

    # Efficiency
    p = tf.add_paragraph()
    p.text = "⚡ Operational Efficiency"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(20)
    p.space_after = Pt(10)

    add_bullet(tf, "Automated barcode scanning reduces manual data entry", 16)
    add_bullet(tf, "Centralized data access from any workstation", 16)
    add_bullet(tf, "Instant traceability queries (seconds vs. hours)", 16)

    # Compliance
    p = tf.add_paragraph()
    p.text = "📋 Compliance & Traceability"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(20)
    p.space_after = Pt(10)

    add_bullet(tf, "Complete audit trail for regulatory requirements", 16)
    add_bullet(tf, "Rapid response to quality issues", 16)
    add_bullet(tf, "Data-driven decision making", 16)

    # ==========================================
    # Slide 10: Investment Summary
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Investment Summary")

    # Investment table
    rows, cols = 9, 3
    left = Inches(2)
    top = Inches(2.2)
    width = Inches(6)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(1.5)

    # Header
    headers = ['Item', 'Quantity', 'Amount (KRW)']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    investment_data = [
        ['Desktop PC (Server)', '1', '1,090,000'],
        ['Fanless Industrial PC', '2', '1,720,000'],
        ['Monitor 24"', '1', '155,000'],
        ['Touch Monitor 24"', '1', '370,000'],
        ['Label Printer', '2', '1,050,000'],
        ['Barcode Scanner', '2', '760,000'],
        ['', '', ''],
        ['Total Investment', '9 units', '5,145,000']
    ]

    for row_idx, row_data in enumerate(investment_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            if row_idx == 8:  # Total row
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(16)

    # ==========================================
    # Save presentation
    # ==========================================
    filename = 'Powder_Inspection_System_Presentation.pptx'
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"📄 Total slides: {len(prs.slides)}")
    return filename


def add_slide_title(slide, title_text):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    # Add line under title
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(0.5), Inches(1.4),
        Inches(9), Inches(0)
    )
    line.line.color.rgb = RGBColor(103, 126, 234)
    line.line.width = Pt(2)


def add_bullet(text_frame, text, font_size=14):
    """Add bullet point to text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = 0
    p.font.size = Pt(font_size)
    p.space_after = Pt(8)


def add_architecture_box(slide, x, y, w, h, title, content):
    """Add architecture component box"""
    # Box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 235, 250)
    shape.line.color.rgb = RGBColor(103, 126, 234)
    shape.line.width = Pt(2)

    # Text
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(8)


def add_workflow_box(slide, x, y, w, h, text):
    """Add workflow step box"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 235, 250)
    shape.line.color.rgb = RGBColor(103, 126, 234)
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, x1, y1, x2, y2):
    """Add arrow connector"""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(103, 126, 234)
    connector.line.width = Pt(2)


if __name__ == '__main__':
    print("🎨 Generating Powder Inspection Management System Presentation...")
    print("=" * 60)
    create_presentation()
    print("=" * 60)
    print("✨ Done! You can now open the presentation file.")
