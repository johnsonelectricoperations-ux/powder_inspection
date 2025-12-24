#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
분말 검사 관리 시스템 - 프레젠테이션 생성기 (한글 버전)
한영 혼용으로 소프트웨어 및 하드웨어 내용 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """분말 검사 관리 시스템 프레젠테이션 생성"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # 슬라이드 1: 표지
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(103, 126, 234)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "분말 검사 관리 시스템"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Powder Inspection Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 설명
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "품질 관리 및 추적성 솔루션"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(20)
    desc_para.font.color.rgb = RGBColor(255, 255, 255)
    desc_para.alignment = PP_ALIGN.CENTER

    # 날짜
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "2024년 12월"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER

    # ==========================================
    # 슬라이드 2: 시스템 개요
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "시스템 개요")

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # 목적
    p = tf.paragraphs[0]
    p.text = "도입 목적"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(tf, "분말 검사 및 배합 작업의 종합적인 품질 관리 시스템 구축", 18)
    add_bullet(tf, "원재료부터 완제품까지 완전한 추적성 (Traceability) 확보", 18)
    add_bullet(tf, "바코드/QR 코드 기반 자동화된 LOT 추적 관리", 18)

    # 주요 특징
    p = tf.add_paragraph()
    p.text = "주요 특징"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(24)
    p.space_after = Pt(12)

    add_bullet(tf, "웹 기반 애플리케이션 (어디서나 접속 가능)", 18)
    add_bullet(tf, "실시간 데이터 동기화", 18)
    add_bullet(tf, "바코드 라벨 출력 및 스캔 시스템", 18)
    add_bullet(tf, "품질 검사 데이터 관리", 18)
    add_bullet(tf, "완전한 전/후방 추적성", 18)

    # ==========================================
    # 슬라이드 3: 시스템 아키텍처
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "시스템 아키텍처")

    # 서버
    add_architecture_box(slide, 4, 2.2, 2, 1.2, "서버 PC (A)",
                         "• Flask 웹 애플리케이션\n• SQLite 데이터베이스\n• 중앙 데이터 관리")

    # 네트워크
    network_box = slide.shapes.add_textbox(Inches(4.5), Inches(3.7), Inches(1), Inches(0.4))
    nf = network_box.text_frame
    nf.text = "네트워크 (LAN)"
    nf.paragraphs[0].font.size = Pt(14)
    nf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 클라이언트 1
    add_architecture_box(slide, 1, 4.5, 2, 1.2, "클라이언트 PC 1",
                         "• 수입검사 스테이션\n• 바코드 스캐너\n• 라벨 프린터")

    # 클라이언트 2
    add_architecture_box(slide, 7, 4.5, 2, 1.2, "클라이언트 PC 2",
                         "• 배합작업 스테이션\n• 바코드 스캐너\n• 라벨 프린터\n• 터치 모니터")

    # ==========================================
    # 슬라이드 4: 소프트웨어 기술 스택
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "소프트웨어 구성")

    # 왼쪽: 기술 스택
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.2), Inches(4.5))
    left_tf = left_box.text_frame

    p = left_tf.paragraphs[0]
    p.text = "기술 스택 (Technology Stack)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(left_tf, "백엔드: Python 3.x + Flask", 16)
    add_bullet(left_tf, "데이터베이스: SQLite (WAL 모드)", 16)
    add_bullet(left_tf, "프론트엔드: HTML5 + CSS3 + JavaScript", 16)
    add_bullet(left_tf, "바코드: JsBarcode 라이브러리", 16)
    add_bullet(left_tf, "통신: RESTful API", 16)

    p = left_tf.add_paragraph()
    p.text = "\n시스템 아키텍처"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)
    p.space_after = Pt(12)

    add_bullet(left_tf, "서버-클라이언트 구조", 16)
    add_bullet(left_tf, "웹 브라우저 기반 접근", 16)
    add_bullet(left_tf, "크로스 플랫폼 지원", 16)

    # 오른쪽: 소프트웨어 특징
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.2), Inches(4.5))
    right_tf = right_box.text_frame

    p = right_tf.paragraphs[0]
    p.text = "소프트웨어 특징"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(right_tf, "설치 불필요 (웹 브라우저만 사용)", 16)
    add_bullet(right_tf, "동시 다중 사용자 지원", 16)
    add_bullet(right_tf, "반응형 웹 디자인", 16)
    add_bullet(right_tf, "실시간 데이터 동기화", 16)
    add_bullet(right_tf, "직관적인 사용자 인터페이스", 16)

    p = right_tf.add_paragraph()
    p.text = "\n데이터 관리"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)
    p.space_after = Pt(12)

    add_bullet(right_tf, "자동 백업 및 복구", 16)
    add_bullet(right_tf, "검사 결과 영구 보존", 16)
    add_bullet(right_tf, "완전한 감사 추적 (Audit Trail)", 16)
    add_bullet(right_tf, "데이터 무결성 보장", 16)

    # ==========================================
    # 슬라이드 5: 주요 기능
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "주요 기능")

    # 왼쪽
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(5))
    left_tf = left_box.text_frame

    p = left_tf.paragraphs[0]
    p.text = "1. 수입분말 검사"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    add_bullet(left_tf, "원재료 품질 검사", 14)
    add_bullet(left_tf, "LOT 번호 등록", 14)
    add_bullet(left_tf, "다항목 검사 결과 관리", 14)

    p = left_tf.add_paragraph()
    p.text = "2. 배합작업 계획 등록"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(left_tf, "배합 규격서 관리", 14)
    add_bullet(left_tf, "작업 지시서 생성", 14)
    add_bullet(left_tf, "생산 일정 관리", 14)

    p = left_tf.add_paragraph()
    p.text = "3. 원재료 투입 추적"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(left_tf, "바코드 스캔 LOT 입력", 14)
    add_bullet(left_tf, "중량 측정 기록", 14)
    add_bullet(left_tf, "이종분말 혼입 감지", 14)

    # 오른쪽
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
    right_tf = right_box.text_frame

    p = right_tf.paragraphs[0]
    p.text = "4. 검사 결과 조회"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    add_bullet(right_tf, "종합 검사 결과 저장", 14)
    add_bullet(right_tf, "합격/불합격 판정", 14)
    add_bullet(right_tf, "이력 데이터 관리", 14)

    p = right_tf.add_paragraph()
    p.text = "5. 추적성 조회"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(right_tf, "후방 추적: 완제품 → 원재료", 14)
    add_bullet(right_tf, "전방 추적: 원재료 → 제품", 14)
    add_bullet(right_tf, "완전한 감사 추적", 14)

    p = right_tf.add_paragraph()
    p.text = "6. 바코드/라벨 시스템"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(18)

    add_bullet(right_tf, "자동 라벨 생성", 14)
    add_bullet(right_tf, "QR/바코드 출력", 14)
    add_bullet(right_tf, "무선 스캐너 통합", 14)

    # ==========================================
    # 슬라이드 6: 사용자 인터페이스
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "사용자 인터페이스 - 메뉴 구조")

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = content_box.text_frame

    menu_items = [
        "1. 대시보드 - 시스템 현황 및 진행중 검사 확인",
        "2. 수입분말 검사 - 원재료 품질 검사",
        "3. 배합작업계획 등록 - 배합 규격서 및 작업 지시서 관리",
        "4. 배합작업 - 원재료 투입 및 생산 실행",
        "5. 배합작업 현황 조회 - 실시간 생산 모니터링",
        "6. 배합분말 검사 - 완제품 품질 검사",
        "7. 검사결과 조회 - 검사 이력 데이터 조회",
        "8. REWORK - 품질 이슈 처리 (향후 구현)",
        "9. 추적성 조회 - LOT 추적 및 감사 추적",
        "10. 관리자 모드 - 시스템 설정 및 사용자 관리"
    ]

    for item in menu_items:
        add_bullet(tf, item, 16)

    # ==========================================
    # 슬라이드 7: 업무 프로세스
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "업무 프로세스")

    # 워크플로우 단계
    steps = [
        ("1. 수입검사", 1, 2.2),
        ("2. 배합규격\n설정", 3, 2.2),
        ("3. 작업지시\n생성", 5, 2.2),
        ("4. 원재료\n투입", 7, 2.2),
        ("5. 배합작업\n실행", 1, 4.5),
        ("6. 라벨\n출력", 3, 4.5),
        ("7. 품질검사", 5, 4.5),
        ("8. 추적성\n관리", 7, 4.5)
    ]

    for step_text, x, y in steps:
        add_workflow_box(slide, x, y, 1.5, 1, step_text)

    # 화살표 추가
    add_arrow(slide, 2.6, 2.7, 2.9, 2.7)  # 1 -> 2
    add_arrow(slide, 4.6, 2.7, 4.9, 2.7)  # 2 -> 3
    add_arrow(slide, 6.6, 2.7, 6.9, 2.7)  # 3 -> 4
    add_arrow(slide, 7.75, 3.3, 7.75, 4.4)  # 4 -> down
    add_arrow(slide, 7.25, 5, 4.6, 5)  # down -> 7
    add_arrow(slide, 4.1, 5, 2.6, 5)  # 7 -> 6
    add_arrow(slide, 2.1, 5, 1.75, 5)  # 6 -> 5
    add_arrow(slide, 1.2, 5.5, 1.2, 6.5)  # 5 -> down
    add_arrow(slide, 1.7, 6.7, 6.9, 6.7)  # across to 8

    # ==========================================
    # 슬라이드 8: 하드웨어 구성
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "하드웨어 구성")

    # 표 생성
    rows, cols = 8, 4
    left = Inches(0.8)
    top = Inches(2)
    width = Inches(8.4)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 컬럼 너비
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.5)

    # 헤더
    headers = ['장비명', '수량', '모델', '용도']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 데이터
    equipment_data = [
        ['데스크톱 PC (서버)', '1', 'SBD71223', '중앙 서버 및 데이터베이스'],
        ['팬리스 산업용 PC', '2', 'HDL-BOX90C', '검사 및 배합 스테이션'],
        ['모니터 24"', '1', 'LG 24MR400W', '서버 디스플레이'],
        ['터치 모니터 24"', '1', 'LG 24MR400', '배합작업 스테이션'],
        ['라벨 프린터', '2', 'BIXOLON SLP-TX420/400', '바코드 라벨 출력'],
        ['바코드 스캐너', '2', 'Zebra DS2278', 'LOT 번호 스캔'],
        ['총 투자 금액', '9대', '', '5,145,000원']
    ]

    for row_idx, row_data in enumerate(equipment_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if row_idx == 7:  # 합계 행
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
                cell.text_frame.paragraphs[0].font.bold = True

    # ==========================================
    # 슬라이드 9: 바코드 및 추적성 시스템
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "바코드 및 추적성 시스템")

    # 왼쪽: 바코드 시스템
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.2), Inches(4.5))
    left_tf = left_box.text_frame

    p = left_tf.paragraphs[0]
    p.text = "바코드 라벨 시스템"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(left_tf, "형식: CODE128 / QR 코드", 14)
    add_bullet(left_tf, "내용: 제품명, LOT번호, 팩번호, 중량, 날짜", 14)
    add_bullet(left_tf, "크기: 100mm × 100mm 열전사 라벨", 14)
    add_bullet(left_tf, "프린터: BIXOLON 열전사 프린터", 14)
    add_bullet(left_tf, "스캐너: Zebra DS2278 무선 2D", 14)

    p = left_tf.add_paragraph()
    p.text = "\n바코드 예시:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(18)

    p = left_tf.add_paragraph()
    p.text = "PN:High-Strength-Mix|LOT:20241224-001|1/5"
    p.font.size = Pt(11)
    p.font.name = 'Courier New'

    # 오른쪽: 추적성 기능
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.2), Inches(4.5))
    right_tf = right_box.text_frame

    p = right_tf.paragraphs[0]
    p.text = "추적성 기능"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(12)

    add_bullet(right_tf, "후방 추적성 (Backward):\n완제품 → 원재료 추적", 14)
    add_bullet(right_tf, "전방 추적성 (Forward):\n원재료 → 사용된 모든 제품 추적", 14)
    add_bullet(right_tf, "타임스탬프 포함 완전한 감사 추적", 14)
    add_bullet(right_tf, "품질 검사 결과 연계", 14)
    add_bullet(right_tf, "작업자 및 장비 추적", 14)

    p = right_tf.add_paragraph()
    p.text = "\n활용 효과:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(18)

    add_bullet(right_tf, "품질 이슈 신속한 원인 파악", 14)
    add_bullet(right_tf, "규제 준수 (Compliance)", 14)
    add_bullet(right_tf, "리콜 대응 능력 확보", 14)

    # ==========================================
    # 슬라이드 10: 기대 효과
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "기대 효과")

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = content_box.text_frame

    # 품질 개선
    p = tf.paragraphs[0]
    p.text = "🎯 품질 개선"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_after = Pt(10)

    add_bullet(tf, "체계적인 품질 검사 프로세스 확립", 16)
    add_bullet(tf, "디지털 기록 관리로 종이 문서 오류 제거", 16)
    add_bullet(tf, "실시간 품질 모니터링", 16)

    # 효율성
    p = tf.add_paragraph()
    p.text = "⚡ 운영 효율성 향상"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(20)
    p.space_after = Pt(10)

    add_bullet(tf, "바코드 스캔으로 수작업 데이터 입력 최소화", 16)
    add_bullet(tf, "어느 워크스테이션에서나 중앙 데이터 접근", 16)
    add_bullet(tf, "즉시 추적성 조회 (수 시간 → 수 초)", 16)

    # 규정 준수
    p = tf.add_paragraph()
    p.text = "📋 규정 준수 및 추적성"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.space_before = Pt(20)
    p.space_after = Pt(10)

    add_bullet(tf, "규제 요구사항 대응 완전한 감사 추적", 16)
    add_bullet(tf, "품질 이슈 발생 시 신속한 대응", 16)
    add_bullet(tf, "데이터 기반 의사 결정 지원", 16)

    # ==========================================
    # 슬라이드 11: 투자 요약
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "투자 요약")

    # 투자 내역 표
    rows, cols = 9, 3
    left = Inches(2)
    top = Inches(2.2)
    width = Inches(6)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 컬럼 너비
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(1.5)

    # 헤더
    headers = ['항목', '수량', '금액 (원)']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 데이터
    investment_data = [
        ['데스크톱 PC (서버)', '1', '1,090,000'],
        ['팬리스 산업용 PC', '2', '1,720,000'],
        ['모니터 24"', '1', '155,000'],
        ['터치 모니터 24"', '1', '370,000'],
        ['라벨 프린터', '2', '1,050,000'],
        ['바코드 스캐너', '2', '760,000'],
        ['', '', ''],
        ['총 투자 금액', '9대', '5,145,000']
    ]

    for row_idx, row_data in enumerate(investment_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            if row_idx == 8:  # 합계 행
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(103, 126, 234)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(16)

    # ==========================================
    # 저장
    # ==========================================
    filename = '분말검사시스템_보고자료.pptx'
    prs.save(filename)
    print(f"✅ 프레젠테이션 생성 완료: {filename}")
    print(f"📄 총 슬라이드 수: {len(prs.slides)}개")
    return filename


def add_slide_title(slide, title_text):
    """슬라이드 제목 추가"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)

    # 제목 밑줄
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(0.5), Inches(1.4),
        Inches(9), Inches(0)
    )
    line.line.color.rgb = RGBColor(103, 126, 234)
    line.line.width = Pt(2)


def add_bullet(text_frame, text, font_size=14):
    """글머리 기호 추가"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = 0
    p.font.size = Pt(font_size)
    p.space_after = Pt(8)


def add_architecture_box(slide, x, y, w, h, title, content):
    """아키텍처 구성요소 박스 추가"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 235, 250)
    shape.line.color.rgb = RGBColor(103, 126, 234)
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(103, 126, 234)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(8)


def add_workflow_box(slide, x, y, w, h, text):
    """워크플로우 단계 박스 추가"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 235, 250)
    shape.line.color.rgb = RGBColor(103, 126, 234)
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, x1, y1, x2, y2):
    """화살표 연결선 추가"""
    from pptx.enum.shapes import MSO_CONNECTOR

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(103, 126, 234)
    connector.line.width = Pt(2)


if __name__ == '__main__':
    print("🎨 분말 검사 관리 시스템 프레젠테이션 생성 중...")
    print("=" * 60)
    create_presentation()
    print("=" * 60)
    print("✨ 완료! 프레젠테이션 파일을 열어보세요.")
